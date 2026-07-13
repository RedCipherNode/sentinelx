from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from reportlab.platypus import SimpleDocTemplate, Paragraph
from reportlab.lib.styles import getSampleStyleSheet
from odf.opendocument import OpenDocumentText
from odf.text import P
from pathlib import Path
from PIL import Image, ImageDraw
import shutil
import zipfile
import json
import yaml
import hashlib
import mimetypes


ROOT = Path(__file__).resolve().parent.parent
OUTPUT = ROOT / "output"
GENERATED = OUTPUT / "generated"
ZIP_FILE = OUTPUT / "sentinelx-testdata.zip"


# ==========================================================
# Utilities
# ==========================================================


def reset_output():
    print("[*] Resetting output directory...")

    if GENERATED.exists():
        shutil.rmtree(GENERATED)

    GENERATED.mkdir(parents=True, exist_ok=True)

    if ZIP_FILE.exists():
        ZIP_FILE.unlink()


def ensure_dir(path: Path):
    path.mkdir(parents=True, exist_ok=True)


def write_text(path: Path, content: str, encoding="utf-8"):
    ensure_dir(path.parent)

    with open(path, "w", encoding=encoding) as f:
        f.write(content)


def lorem(paragraphs=5):
    p = (
        "Lorem ipsum dolor sit amet, consectetur adipiscing elit. "
        "Sed do eiusmod tempor incididunt ut labore et dolore magna aliqua.\n\n"
    )

    return p * paragraphs


# ==========================================================
# Text Files
# ==========================================================


def generate_text():
    print("[*] Generating text files...")

    path = GENERATED / "files" / "text"

    write_text(path / "empty.txt", "")
    write_text(path / "lorem.txt", lorem(10))
    write_text(path / "unicode.txt", "こんにちは\n안녕하세요\n你好\n😀\n")
    write_text(path / "utf8.txt", lorem(5))
    write_text(path / "long.txt", lorem(1200))

    ensure_dir(path)

    with open(path / "utf16.txt", "w", encoding="utf-16") as f:
        f.write(lorem(5))


# ==========================================================
# Documents
# ==========================================================


def generate_documents():
    print("[*] Generating documents...")

    path = GENERATED / "files" / "documents"
    ensure_dir(path)

    title = "SentinelX Test Document"
    body = lorem(15)

    # --------------------------------------------------
    # PDF
    # --------------------------------------------------

    pdf = path / "sample.pdf"

    styles = getSampleStyleSheet()

    doc = SimpleDocTemplate(str(pdf))

    story = [
        Paragraph(f"<b>{title}</b>", styles["Heading1"]),
        Paragraph(body.replace("\n", "<br/>"), styles["BodyText"]),
    ]

    doc.build(story)

    # --------------------------------------------------
    # DOCX
    # --------------------------------------------------

    docx = Document()

    docx.add_heading(title, level=1)
    docx.add_paragraph(body)

    docx.save(path / "sample.docx")

    # --------------------------------------------------
    # XLSX
    # --------------------------------------------------

    wb = Workbook()

    ws = wb.active
    ws.title = "Sheet1"

    ws["A1"] = title

    for i in range(2, 22):
        ws[f"A{i}"] = f"Lorem Ipsum Row {i - 1}"

    wb.save(path / "sample.xlsx")

    # --------------------------------------------------
    # PPTX
    # --------------------------------------------------

    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[5])

    slide.shapes.title.text = title

    textbox = slide.shapes.add_textbox(
        left=500000,
        top=1000000,
        width=7000000,
        height=3000000,
    )

    textbox.text_frame.text = body

    prs.save(path / "sample.pptx")

    # --------------------------------------------------
    # ODT
    # --------------------------------------------------

    odt = OpenDocumentText()

    odt.text.addElement(P(text=title))
    odt.text.addElement(P(text=body))

    odt.save(str(path / "sample.odt"))

    # --------------------------------------------------
    # RTF
    # --------------------------------------------------

    rtf = r"""{\rtf1\ansi

\b SentinelX Test Document\b0

\par

Lorem ipsum dolor sit amet.

\par

Lorem ipsum dolor sit amet.

\par

Lorem ipsum dolor sit amet.

}
"""

    write_text(path / "sample.rtf", rtf)


# ==========================================================
# Images
# ==========================================================


def generate_images():
    print("[*] Generating images...")

    path = GENERATED / "files" / "images"
    ensure_dir(path)

    for ext in [
        "png",
        "jpg",
        "gif",
        "bmp",
        "webp",
    ]:
        img = Image.new(
            "RGB",
            (800, 600),
            (52, 120, 246),
        )

        draw = ImageDraw.Draw(img)

        draw.text(
            (220, 290),
            "SentinelX Test Image",
            fill="white",
        )

        img.save(path / f"sample.{ext}")


# ==========================================================
# Executables
# ==========================================================


def generate_executables():
    print("[*] Generating fake executables...")

    path = GENERATED / "files" / "executables"
    ensure_dir(path)

    content = """This is NOT a real executable.

This file exists only for extension testing.

SentinelX Test Dataset.
"""

    for filename in [
        "fake.exe",
        "fake.dll",
        "fake.sys",
        "renamed_pdf.exe",
        "renamed_image.exe",
    ]:
        write_text(path / filename, content)


# ==========================================================
# Scripts
# ==========================================================


def generate_scripts():
    print("[*] Generating scripts...")

    path = GENERATED / "files" / "scripts"
    ensure_dir(path)

    scripts = {
        "hello.ps1": 'Write-Output "Hello SentinelX"\n',
        "hello.bat": "@echo off\necho Hello SentinelX\n",
        "hello.cmd": "@echo off\necho Hello SentinelX\n",
        "hello.sh": "#!/bin/sh\necho Hello SentinelX\n",
        "hello.py": 'print("Hello SentinelX")\n',
        "hello.js": 'console.log("Hello SentinelX");\n',
    }

    for filename, content in scripts.items():
        write_text(path / filename, content)


# ==========================================================
# Miscellaneous
# ==========================================================


def generate_misc():
    print("[*] Generating miscellaneous files...")

    path = GENERATED / "files" / "miscellaneous"
    ensure_dir(path)

    data = {
        "project": "SentinelX",
        "version": "1.0.0",
        "safe": True,
        "description": "Generated test data",
    }

    # JSON
    with open(path / "sample.json", "w", encoding="utf-8") as f:
        json.dump(data, f, indent=4)

    # XML
    write_text(
        path / "sample.xml",
        """<?xml version="1.0"?>
<project>
    <name>SentinelX</name>
    <version>1.0.0</version>
</project>
""",
    )

    # YAML
    with open(path / "sample.yaml", "w", encoding="utf-8") as f:
        yaml.dump(data, f, sort_keys=False)

    # TOML
    write_text(
        path / "sample.toml",
        """
project = "SentinelX"
version = "1.0.0"
safe = true
""".strip(),
    )

    # INI
    write_text(
        path / "sample.ini",
        """
[project]
name=SentinelX
version=1.0.0
""".strip(),
    )

    # CSV
    write_text(
        path / "sample.csv",
        """id,name,role
1,Alice,Developer
2,Bob,Tester
3,Charlie,User
""",
    )

    # LOG
    write_text(
        path / "sample.log",
        """INFO Starting SentinelX
INFO Loading configuration
INFO Done
""",
    )


# ==========================================================
# Projects
# ==========================================================


def generate_projects():
    print("[*] Generating projects...")

    root = GENERATED / "projects"

    # ======================================================
    # Node.js
    # ======================================================

    node = root / "node"
    ensure_dir(node)

    write_text(
        node / "package.json",
        """{
  "name": "sentinelx-test-node",
  "version": "1.0.0",
  "description": "SentinelX test project",
  "main": "index.js",
  "license": "MIT"
}
""",
    )

    write_text(
        node / "README.md",
        "# SentinelX Node Project\n\nLorem ipsum dolor sit amet.\n",
    )

    write_text(
        node / "index.js",
        'console.log("Hello SentinelX");\n',
    )

    # ======================================================
    # Python
    # ======================================================

    py = root / "python"
    ensure_dir(py)

    write_text(
        py / "requirements.txt",
        "requests\n",
    )

    write_text(
        py / "main.py",
        'print("Hello SentinelX")\n',
    )

    # ======================================================
    # Rust
    # ======================================================

    rust = root / "rust"
    ensure_dir(rust / "src")

    write_text(
        rust / "Cargo.toml",
        """[package]
name = "sentinelx-test"
version = "0.1.0"
edition = "2024"
""",
    )

    write_text(
        rust / "src" / "main.rs",
        """fn main() {
    println!("Hello SentinelX");
}
""",
    )

    # ======================================================
    # Go
    # ======================================================

    go = root / "go"
    ensure_dir(go)

    write_text(
        go / "go.mod",
        """module sentinelx-test

go 1.22
""",
    )

    write_text(
        go / "main.go",
        """package main

import "fmt"

func main() {
    fmt.Println("Hello SentinelX")
}
""",
    )

    # ======================================================
    # C
    # ======================================================

    c = root / "c"
    ensure_dir(c)

    write_text(
        c / "main.c",
        """#include <stdio.h>

int main(void)
{
    printf("Hello SentinelX\\n");
    return 0;
}
""",
    )

    write_text(
        c / "Makefile",
        """all:
\tgcc main.c -o app
""",
    )

    # ======================================================
    # C++
    # ======================================================

    cpp = root / "cpp"
    ensure_dir(cpp)

    write_text(
        cpp / "main.cpp",
        """#include <iostream>

int main()
{
    std::cout << "Hello SentinelX" << std::endl;
    return 0;
}
""",
    )

    write_text(
        cpp / "CMakeLists.txt",
        """cmake_minimum_required(VERSION 3.15)

project(SentinelXTest)

add_executable(app main.cpp)
""",
    )

    # ======================================================
    # Java
    # ======================================================

    java = root / "java"
    ensure_dir(java)

    write_text(
        java / "pom.xml",
        """<project xmlns="http://maven.apache.org/POM/4.0.0">
  <modelVersion>4.0.0</modelVersion>

  <groupId>dev.sentinelx</groupId>
  <artifactId>sample</artifactId>
  <version>1.0.0</version>
</project>
""",
    )

    write_text(
        java / "Main.java",
        """public class Main {

    public static void main(String[] args) {
        System.out.println("Hello SentinelX");
    }

}
""",
    )

    # ======================================================
    # .NET
    # ======================================================

    dotnet = root / "dotnet"
    ensure_dir(dotnet)

    write_text(
        dotnet / "Program.cs",
        """Console.WriteLine("Hello SentinelX");
""",
    )

    write_text(
        dotnet / "sample.csproj",
        """<Project Sdk="Microsoft.NET.Sdk">

  <PropertyGroup>
    <OutputType>Exe</OutputType>
    <TargetFramework>net8.0</TargetFramework>
  </PropertyGroup>

</Project>
""",
    )


# ==========================================================
# URLs
# ==========================================================


def generate_urls():
    print("[*] Generating URL list...")

    path = GENERATED / "urls"
    ensure_dir(path)

    urls = [
        "https://google.com",
        "https://github.com",
        "https://docs.rs",
        "https://crates.io",
        "https://npmjs.com",
        "https://pypi.org",
        "https://wikipedia.org",
        "http://localhost",
        "http://localhost:8080",
        "http://127.0.0.1",
        "http://127.0.0.1:5000",
    ]

    for i in range(1, 41):
        urls.append(f"https://example.com/search?q=sentinelx&page={i}&sort=desc")

    write_text(path / "urls.txt", "\n".join(urls))


# ==========================================================
# Commands
# ==========================================================


def generate_commands():
    print("[*] Generating command examples...")

    path = GENERATED / "commands"
    ensure_dir(path)

    commands = {
        "powershell.txt": [
            "Get-Process",
            "Get-Service",
            "Get-ChildItem",
        ],
        "bash.txt": [
            "ls",
            "pwd",
            "echo Hello",
        ],
        "cmd.txt": [
            "dir",
            "ipconfig",
            "whoami",
        ],
        "npm.txt": [
            "npm install",
            "npm run build",
            "npm publish",
        ],
        "pip.txt": [
            "pip install requests",
            "pip list",
        ],
        "cargo.txt": [
            "cargo build",
            "cargo run",
            "cargo test",
        ],
        "docker.txt": [
            "docker build .",
            "docker run image",
        ],
        "git.txt": [
            "git clone https://github.com/example/repo.git",
            "git status",
            'git commit -m "Example"',
        ],
        "kubectl.txt": [
            "kubectl get pods",
            "kubectl describe pod example",
        ],
    }

    for filename, lines in commands.items():
        write_text(path / filename, "\n".join(lines))


# ==========================================================
# Edge Cases
# ==========================================================


def generate_edge_cases():
    print("[*] Generating edge cases...")

    root = GENERATED / "edge_cases"

    # ======================================================
    # Empty files
    # ======================================================

    empty = root / "empty"
    ensure_dir(empty)

    for filename in [
        "empty.txt",
        "empty.json",
        "empty.pdf",
        "empty.png",
        "empty.zip",
    ]:
        write_text(empty / filename, "")

    # ======================================================
    # Filename edge cases
    # ======================================================

    names = root / "filenames"
    ensure_dir(names)

    files = [
        "file with spaces.txt",
        "UPPERCASE.TXT",
        "MiXeD-CaSe.txt",
        "Üñïçødê.txt",
        "日本語.txt",
        "emoji-😀.txt",
        ".hidden",
        "no_extension",
        ".".join(["many"] * 12) + ".txt",
        ("very-" * 20) + "long-name.txt",
    ]

    for file in files:
        write_text(names / file, "SentinelX Edge Case")

    # ======================================================
    # Directories
    # ======================================================

    ensure_dir(root / "directories" / "empty-folder")

    ensure_dir(root / "directories" / "nested" / "a" / "b" / "c" / "d" / "e")

    ensure_dir(root / "directories" / "folder with spaces")

    # ======================================================
    # Extension mismatch
    # ======================================================

    ext = root / "extensions"
    ensure_dir(ext)

    write_text(ext / "photo.png.txt", "Not a PNG")
    write_text(ext / "invoice.pdf.exe", "Not an executable")
    write_text(ext / "archive.zip.pdf", "Not a PDF")
    write_text(ext / "image.jpg.json", "{}")
    write_text(ext / "config.toml.txt", "name = SentinelX")


# ==========================================================
# Archives
# ==========================================================


def generate_archives():
    print("[ ] TODO Archives")


# ==========================================================
# Manifest
# ==========================================================


def generate_manifest():
    print("[*] Generating manifest...")

    manifest = []

    for file in GENERATED.rglob("*"):
        if not file.is_file():
            continue

        sha256 = hashlib.sha256(file.read_bytes()).hexdigest()

        manifest.append(
            {
                "path": str(file.relative_to(GENERATED)).replace("\\", "/"),
                "size": file.stat().st_size,
                "mime": mimetypes.guess_type(file)[0],
                "sha256": sha256,
            }
        )

    with open(
        GENERATED / "manifest.json",
        "w",
        encoding="utf-8",
    ) as f:
        json.dump(
            manifest,
            f,
            indent=4,
        )


# ==========================================================
# ZIP
# ==========================================================


def create_zip():
    print("[*] Creating ZIP...")

    with zipfile.ZipFile(ZIP_FILE, "w", zipfile.ZIP_DEFLATED) as z:
        for file in GENERATED.rglob("*"):
            if file.is_file():
                z.write(file, file.relative_to(OUTPUT))


# ==========================================================
# Main
# ==========================================================


def main():
    print("====================================")
    print(" SentinelX Test Data Generator")
    print("====================================")

    reset_output()

    generate_text()
    generate_documents()
    generate_images()
    generate_executables()
    generate_scripts()
    generate_misc()
    generate_projects()
    generate_urls()
    generate_commands()
    generate_edge_cases()
    generate_archives()

    generate_manifest()

    create_zip()

    print()
    print("Done.")
    print(f"Dataset : {GENERATED}")
    print(f"ZIP     : {ZIP_FILE}")


if __name__ == "__main__":
    main()
